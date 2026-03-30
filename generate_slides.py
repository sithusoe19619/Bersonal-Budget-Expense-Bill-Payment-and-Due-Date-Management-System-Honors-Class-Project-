from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

BG_DARK    = RGBColor(0x0F, 0x0F, 0x1A)
BG_CARD    = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT     = RGBColor(0x6C, 0x63, 0xFF)
ACCENT2    = RGBColor(0x00, 0xD4, 0xFF)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xDD)
DIM        = RGBColor(0x88, 0x88, 0xAA)

UML_IMAGE = "/Users/marshall/Desktop/BMCC All Semesters/SPRING26/CSC 331H/Honors Project/finance-manager/umlInitial.png"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def set_bg(slide, color=BG_DARK):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def box(slide, text, l, t, w, h, size=18, bold=False,
        color=WHITE, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Segoe UI" if not bold else "Segoe UI Semibold"
    return tb

def rect(slide, l, t, w, h, fill_color, alpha=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    s.line.fill.background()
    return s

def gradient_bar(slide):
    rect(slide, 0, 0, 13.33, 0.06, ACCENT)

def make_title_slide():
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    # Left accent column
    rect(slide, 0, 0, 0.08, 7.5, ACCENT)
    # Decorative circle top right
    circ = slide.shapes.add_shape(9, Inches(10.5), Inches(-1.2), Inches(4.5), Inches(4.5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x3E)
    circ.line.fill.background()
    circ2 = slide.shapes.add_shape(9, Inches(11.2), Inches(4.8), Inches(2.8), Inches(2.8))
    circ2.fill.solid()
    circ2.fill.fore_color.rgb = RGBColor(0x16, 0x16, 0x32)
    circ2.line.fill.background()
    # Tag line
    box(slide, "CSC 331H  ·  Honors Project  ·  Spring 2026",
        0.5, 1.5, 10, 0.4, size=13, color=ACCENT2, italic=True)
    # Main title
    box(slide, "Personal Budget &\nExpense Management\nSystem",
        0.5, 2.0, 10, 2.6, size=46, bold=True, color=WHITE)
    # Subtitle
    box(slide, "Design Phase Presentation",
        0.5, 4.75, 9, 0.5, size=20, color=LIGHT_GRAY)
    # Bottom divider
    rect(slide, 0.5, 5.55, 6, 0.04, ACCENT)
    box(slide, "BMCC  ·  Department of Computer Information Systems",
        0.5, 5.75, 9, 0.4, size=13, color=DIM)
    return slide

def make_section_slide(section_num, section_title, description):
    """Interstitial section divider slide."""
    slide = prs.slides.add_slide(blank)
    set_bg(slide, BG_CARD)
    rect(slide, 0, 0, 0.08, 7.5, ACCENT2)
    box(slide, f"0{section_num}", 0.5, 1.8, 3, 1.8, size=96, bold=True,
        color=RGBColor(0x2A, 0x2A, 0x4A))
    box(slide, section_title, 0.5, 3.2, 12, 0.8, size=34, bold=True, color=WHITE)
    box(slide, description, 0.5, 4.15, 11, 0.6, size=17, color=LIGHT_GRAY)
    return slide

def make_content_slide(title, subtitle, paragraphs):
    """
    paragraphs: list of (heading, body_text) tuples.
    heading can be None for plain body paragraphs.
    """
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    gradient_bar(slide)
    # Title
    box(slide, title, 0.45, 0.18, 12.4, 0.65, size=28, bold=True, color=WHITE)
    # Subtitle rule
    if subtitle:
        box(slide, subtitle, 0.45, 0.85, 12.4, 0.35, size=14,
            color=ACCENT2, italic=True)
    rect(slide, 0.45, 1.28, 12.4, 0.03, ACCENT)

    top = 1.42
    for heading, body in paragraphs:
        if heading:
            box(slide, heading, 0.5, top, 12.2, 0.38,
                size=16, bold=True, color=ACCENT2)
            top += 0.38
        if body:
            box(slide, body, 0.6, top, 12.1, 0.7,
                size=15, color=LIGHT_GRAY, wrap=True)
            top += 0.68
    return slide

def make_uml_slide():
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    gradient_bar(slide)
    box(slide, "UML Class Diagram — System Architecture",
        0.45, 0.12, 12.4, 0.6, size=28, bold=True, color=WHITE)
    box(slide, "7 components  ·  9 relationships  ·  3 composition  ·  3 dependency  ·  3 structural",
        0.45, 0.75, 12.4, 0.35, size=13, color=ACCENT2, italic=True)
    rect(slide, 0.45, 1.15, 12.4, 0.03, ACCENT)
    slide.shapes.add_picture(UML_IMAGE,
        Inches(0.35), Inches(1.25), Inches(12.63), Inches(6.1))
    return slide

# ═══════════════════════════════════════════════════════
#  SLIDE 1 — Title
# ═══════════════════════════════════════════════════════
make_title_slide()

# ═══════════════════════════════════════════════════════
#  SLIDE 2 — Project Overview
# ═══════════════════════════════════════════════════════
make_content_slide(
    "Project Overview",
    "A C++ personal finance management system built entirely from scratch — no STL containers.",
    [
        ("What This System Does",
         "This system allows users to log and categorize expenses, track bill due dates, and query "
         "their spending history over any date range. It enforces per-category budget limits and "
         "immediately notifies the user when a limit is exceeded."),
        ("Three Core Actions",
         "Action 1 — Log an expense and check if it exceeds a category budget.  "
         "Action 2 — View the next upcoming bill and mark it as paid.  "
         "Action 3 — Query spending history by date range."),
        ("System Composition",
         "The system is composed of 7 components: Expense, Bill, CategoryInfo, HashMap, MinHeap, "
         "BST, and BudgetManager. All data structures are implemented from scratch. "
         "8 header files define the full interface, reviewed and finalized before implementation."),
    ]
)

# ═══════════════════════════════════════════════════════
#  SLIDE 3 — UML Class Diagram
# ═══════════════════════════════════════════════════════
make_uml_slide()

# ═══════════════════════════════════════════════════════
#  SLIDE 4 — Data Structure Design
# ═══════════════════════════════════════════════════════
make_content_slide(
    "Data Structure Design",
    "Each data structure was chosen because it is the optimal fit for the operation it must support.",
    [
        ("HashMap for Budget Tracking  —  O(1) lookup",
         "Every time an expense is logged, the system must look up the category, update its running "
         "total, and compare it against the budget limit. This happens on every single log operation. "
         "A HashMap provides O(1) average-case lookup regardless of how many categories exist, "
         "making it the only structure fast enough for this use case. It uses open-addressing with "
         "linear probing and an isDeleted tombstone flag to safely support removal."),
        ("MinHeap for Bills  —  O(1) access to earliest due date",
         "Action 2 always needs the bill with the earliest due date. A MinHeap keeps that bill at "
         "the root at all times, delivering it in O(1). Insertion and removal cost O(log n). "
         "A plain array would require scanning all bills — O(n) — on every request."),
        ("BST for Expenses  —  O(log n + k) range queries",
         "Action 3 retrieves all expenses within a date range. A BST stores expenses sorted by "
         "date, so the system descends to the start date in O(log n) and collects all matching "
         "nodes without visiting the rest. A flat array would scan every expense — O(n) — "
         "on every query, degrading as history grows."),
    ]
)

# ═══════════════════════════════════════════════════════
#  SLIDE 5 — Module Specification & Interface Design
# ═══════════════════════════════════════════════════════
make_content_slide(
    "Module Specification & Interface Design",
    "Every component exposes a well-defined public interface. All internal complexity is encapsulated.",
    [
        ("BudgetManager — The Single Entry Point",
         "BudgetManager is the only class the outside world interacts with directly. It owns BST, "
         "MinHeap, and HashMap as private members and exposes a clean public interface that hides "
         "all data structure complexity. No external code ever calls BST::insert() or "
         "HashMap::get() directly — every operation goes through BudgetManager."),
        ("Component Interfaces",
         "BudgetManager delegates to BST for expense storage and range queries, to MinHeap for "
         "bill queue management, and to HashMap for category budget lookups. HashMap in turn "
         "communicates with CategoryInfo to update totals and check limits. "
         "Date is a shared value type passed by const reference across all interfaces."),
        ("Design Safety",
         "The Rule of Three is applied to HashMap, MinHeap, and BST — all three own raw pointers "
         "and have their copy constructor and copy assignment operator deleted to prevent "
         "shallow copies and double-frees. All constructors zero-initialize numeric fields "
         "to eliminate undefined behavior from the start."),
    ]
)

# ═══════════════════════════════════════════════════════
#  SLIDE 6 — Algorithm Definition
# ═══════════════════════════════════════════════════════
make_content_slide(
    "Algorithm Definition",
    "Precise logic is defined for each core action and each supporting data structure operation.",
    [
        ("Action 1 — Log an Expense and Check Budget",
         "addExpense inserts the expense into the BST, then retrieves or creates the matching "
         "CategoryInfo in the HashMap. It calls categoryInfo.addExpense(amount) to increment "
         "totalSpent, then calls checkBudget, which prints a warning if totalSpent exceeds "
         "budgetLimit."),
        ("Action 2 — View Next Bill and Mark as Paid",
         "hasPendingBills checks MinHeap::isEmpty() before any access. getNextBill returns "
         "heap[0] — the root — in O(1). markBillPaid delegates to markPaidByName, which scans "
         "the heap array linearly, sets isPaid = true and paidOn = paymentDate in place. "
         "No re-heapify is needed because the due date does not change."),
        ("Action 3 — Query Spending History by Date Range",
         "rangeQuery calls rangeHelper recursively. If a node's date is before the start, only "
         "the right subtree is visited. If it is after the end, only the left subtree is visited. "
         "If it falls within the range, it is added to the result and both subtrees are explored. "
         "This achieves O(log n + k) where k is the number of results returned."),
    ]
)

# ═══════════════════════════════════════════════════════
#  SLIDE 7 — Data Flow Diagrams
# ═══════════════════════════════════════════════════════
make_content_slide(
    "Data Flow Diagrams",
    "These diagrams trace the path of data from user input to system output for each core action.",
    [
        ("Action 1",
         "User provides category, amount, date, and description  →  addExpense inserts into BST  "
         "→  HashMap retrieves or creates CategoryInfo  →  totalSpent is updated  →  "
         "checkBudget compares totals  →  if over budget, a warning with the excess amount is "
         "printed to the user."),
        ("Action 2",
         "User requests next bill  →  hasPendingBills checks isEmpty  →  if bills exist, "
         "getNextBill calls peek and returns heap[0]  →  bill details are displayed  →  user "
         "confirms payment  →  markBillPaid calls markPaidByName  →  bill is updated in place."),
        ("Action 3",
         "User provides start date and end date  →  getExpensesByRange calls rangeQuery  →  "
         "rangeHelper traverses the BST visiting only relevant nodes  →  matching expenses are "
         "collected into a vector  →  each expense's date, category, amount, and description "
         "are printed to the user."),
    ]
)

# ═══════════════════════════════════════════════════════
#  SLIDE 8 — Unit Test Planning
# ═══════════════════════════════════════════════════════
make_content_slide(
    "Unit Test Planning",
    "30 tests are planned across all 6 modules. Every test maps to a declared method in include/.",
    [
        ("Test Coverage by Module",
         "Date — 5 tests covering constructors, isBefore, operator==, and isBetween.  "
         "CategoryInfo — 5 tests covering initialization, addExpense accumulation, isOverBudget, "
         "and getRemainingBudget.  HashMap — 5 tests covering insert/get round-trip, contains, "
         "remove with tombstone, probe chain integrity, and resize.  MinHeap — 5 tests covering "
         "isEmpty, peek ordering, extractMin, markPaidByName, and unknown-name handling.  "
         "BST — 5 tests covering inOrder sorting, duplicate dates, search returning a vector, "
         "rangeQuery, and empty range.  BudgetManager — 5 integration tests covering all 3 "
         "actions end-to-end plus edge cases."),
        ("Notable Edge Case Tests",
         "Test 14 verifies that removing a HashMap entry does not break the probe chain for "
         "entries that were displaced past it during insertion. Test 22 verifies that two expenses "
         "on the same date are both preserved in the BST using the equal-dates-go-right strategy. "
         "Test 27 verifies that addExpense gracefully auto-creates a CategoryInfo entry when no "
         "budget limit has been set for that category yet."),
    ]
)

# ═══════════════════════════════════════════════════════
#  SLIDE 9 — Traceability Matrix
# ═══════════════════════════════════════════════════════
make_content_slide(
    "Traceability Matrix",
    "Every method in the system traces back to one of the three core user requirements.",
    [
        ("Action 1 — Log an Expense and Check Budget",
         "Fulfilled by: BudgetManager::addExpense, BudgetManager::checkBudget, "
         "BudgetManager::setBudgetLimit, BST::insert, HashMap::get, HashMap::insert, "
         "CategoryInfo::addExpense, CategoryInfo::isOverBudget, and "
         "CategoryInfo::getRemainingBudget."),
        ("Action 2 — View Next Bill and Mark as Paid",
         "Fulfilled by: BudgetManager::hasPendingBills, BudgetManager::getNextBill, "
         "BudgetManager::markBillPaid, MinHeap::isEmpty, MinHeap::peek, "
         "MinHeap::markPaidByName, Bill::markAsPaid, and Bill::display."),
        ("Action 3 — Query Spending History by Date Range",
         "Fulfilled by: BudgetManager::getExpensesByRange, BST::rangeQuery, "
         "BST::rangeHelper, Date::isBefore, Date::isAfter, Date::isBetween, "
         "and Expense::display."),
    ]
)

# ═══════════════════════════════════════════════════════
#  SLIDE 10 — Summary
# ═══════════════════════════════════════════════════════
make_content_slide(
    "Summary",
    "The design phase is complete. All deliverables are documented and the system is ready for implementation.",
    [
        ("What Was Accomplished",
         "A UML class diagram was produced with 7 components and 9 relationships. Eight header "
         "files were written, reviewed across multiple passes, and corrected for uninitialized "
         "fields, missing constructors, pointer safety, and interface gaps. Algorithms were "
         "defined in pseudocode for all three core actions and all key data structure operations."),
        ("Documentation Produced",
         "The reflection.md file contains the complete Detailed Design Document, including the "
         "component brainstorm, data structure justification, design changes, traceability matrix, "
         "algorithm definition, interface design, data flow diagrams, and unit test plan with "
         "30 tests mapped to declared method signatures."),
        ("Next Step — Phase 2: Implementation",
         "All .cpp stub files are in place in src/. The Makefile is configured and ready. "
         "Implementation begins with the foundational components — Date, Expense, Bill, and "
         "CategoryInfo — before moving on to the data structures and BudgetManager."),
    ]
)

# Save
out = "/Users/marshall/Desktop/BMCC All Semesters/SPRING26/CSC 331H/Honors Project/finance-manager/DesignPhase_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Total slides: {len(prs.slides)}")
